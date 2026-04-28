import json
import re
import shutil
import subprocess
from datetime import date
from pathlib import Path

import fitz
import pdfplumber
from docx import Document
from pptx import Presentation
from PIL import Image
import pytesseract

ROOT = Path('/home/smayan/Desktop/QPM')
TARGET_DIRS = [ROOT / 'UNIT 4', ROOT / 'UNIT 5', ROOT / 'UNIT 6']
OUT_JSON = ROOT / 'extracted' / 'unit456_extraction.json'
OUT_MD = ROOT / 'extracted' / 'unit456_extraction.md'
IMG_DIR = ROOT / 'extracted' / 'images'
IMG_DIR.mkdir(parents=True, exist_ok=True)

records = []
TESS_OK = shutil.which('tesseract') is not None


def ocr_image(path: Path) -> str:
    if not TESS_OK:
        return ''
    try:
        txt = pytesseract.image_to_string(Image.open(path))
        txt = re.sub(r'\s+', ' ', txt).strip()
        return txt[:600]
    except Exception:
        return ''


def add_record(path: Path, rec_type: str, data: dict) -> None:
    records.append({'file': str(path.relative_to(ROOT)), 'type': rec_type, **data})


for d in TARGET_DIRS:
    for ppt in d.rglob('*.ppt'):
        outdir = ppt.parent
        cmd = ['soffice', '--headless', '--convert-to', 'pptx', '--outdir', str(outdir), str(ppt)]
        try:
            subprocess.run(cmd, check=True, capture_output=True, text=True)
            converted = ppt.with_suffix('.pptx')
            if converted.exists():
                add_record(ppt, 'conversion', {'status': 'converted', 'to': str(converted.relative_to(ROOT))})
            else:
                add_record(ppt, 'conversion', {'status': 'conversion-command-ran-no-output'})
        except Exception as e:
            add_record(ppt, 'conversion', {'status': 'failed', 'error': str(e)})

for d in TARGET_DIRS:
    for f in sorted(d.rglob('*')):
        if not f.is_file():
            continue
        ext = f.suffix.lower()

        if ext == '.ppt':
            try:
                doc = fitz.open(str(f))
                pages = []
                image_notes = []
                for pno in range(len(doc)):
                    page = doc[pno]
                    txt = page.get_text() or ''
                    txt = re.sub(r'\s+\n', '\n', txt).strip()

                    # If slide has no text layer, OCR a rendered page image.
                    page_img = IMG_DIR / f"{f.stem}_legacy_p{pno + 1}.png"
                    pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                    pix.save(str(page_img))
                    ocr_txt = ocr_image(page_img)

                    pages.append({
                        'page': pno + 1,
                        'text': txt[:6000],
                        'ocr_page': ocr_txt,
                    })
                    image_notes.append({
                        'page': pno + 1,
                        'image': str(page_img.relative_to(ROOT)),
                        'ocr': ocr_txt,
                    })

                add_record(f, 'legacy_ppt_pdflike', {'pages': pages, 'images': image_notes})
            except Exception as e:
                add_record(f, 'legacy_ppt_pdflike', {'error': str(e)})

            continue

        if ext == '.pdf':
            pages = []
            image_notes = []
            try:
                with pdfplumber.open(str(f)) as pdf:
                    for i, p in enumerate(pdf.pages, start=1):
                        t = p.extract_text() or ''
                        t = re.sub(r'\s+\n', '\n', t).strip()
                        pages.append({'page': i, 'text': t[:6000]})
                doc = fitz.open(str(f))
                for pno in range(len(doc)):
                    page = doc[pno]
                    imgs = page.get_images(full=True)
                    for j, img in enumerate(imgs):
                        xref = img[0]
                        base = doc.extract_image(xref)
                        img_bytes = base.get('image')
                        img_ext = base.get('ext', 'png')
                        img_path = IMG_DIR / f"{f.stem}_p{pno + 1}_{j}.{img_ext}"
                        with open(img_path, 'wb') as w:
                            w.write(img_bytes)
                        image_notes.append({
                            'page': pno + 1,
                            'image': str(img_path.relative_to(ROOT)),
                            'ocr': ocr_image(img_path),
                        })
                add_record(f, 'pdf', {'pages': pages, 'images': image_notes})
            except Exception as e:
                add_record(f, 'pdf', {'error': str(e)})

        elif ext == '.docx':
            try:
                doc = Document(str(f))
                paras = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
                tables = []
                for t in doc.tables:
                    rows = []
                    for r in t.rows:
                        rows.append([c.text.strip() for c in r.cells])
                    tables.append(rows)

                images = []
                idx = 0
                for rel in doc.part.rels.values():
                    if 'image' in rel.reltype:
                        idx += 1
                        part = rel.target_part
                        ctype = part.content_type.split('/')[-1]
                        ip = IMG_DIR / f"{f.stem}_docximg{idx}.{ctype}"
                        with open(ip, 'wb') as w:
                            w.write(part.blob)
                        images.append({'image': str(ip.relative_to(ROOT)), 'ocr': ocr_image(ip)})
                add_record(f, 'docx', {'paragraphs': paras[:1500], 'tables': tables[:80], 'images': images})
            except Exception as e:
                add_record(f, 'docx', {'error': str(e)})

        elif ext == '.pptx':
            try:
                prs = Presentation(str(f))
                slides = []
                images = []
                for sno, slide in enumerate(prs.slides, start=1):
                    texts = []
                    for shp in slide.shapes:
                        if hasattr(shp, 'text') and shp.text and shp.text.strip():
                            texts.append(shp.text.strip())
                        if getattr(shp, 'shape_type', None) == 13:
                            try:
                                im = shp.image
                                ip = IMG_DIR / f"{f.stem}_s{sno}_{len(images) + 1}.{im.ext}"
                                with open(ip, 'wb') as w:
                                    w.write(im.blob)
                                images.append({'slide': sno, 'image': str(ip.relative_to(ROOT)), 'ocr': ocr_image(ip)})
                            except Exception:
                                pass

                    notes_text = ''
                    try:
                        ns = slide.notes_slide
                        if ns and ns.notes_text_frame and ns.notes_text_frame.text:
                            notes_text = ns.notes_text_frame.text.strip()
                    except Exception:
                        pass
                    slides.append({'slide': sno, 'text': texts, 'notes': notes_text})
                add_record(f, 'pptx', {'slides': slides, 'images': images})
            except Exception as e:
                add_record(f, 'pptx', {'error': str(e)})

OUT_JSON.write_text(json.dumps(records, indent=2), encoding='utf-8')

lines = []
lines.append(f"# Unit 4-6 Extraction Snapshot ({date.today().isoformat()})")
lines.append(f"- Tesseract available: {TESS_OK}")
for r in records:
    lines.append(f"\n## {r['file']}")
    lines.append(f"- Type: {r['type']}")
    if 'error' in r:
        lines.append(f"- Error: {r['error']}")
        continue
    if r['type'] == 'pdf':
        lines.append(f"- Pages extracted: {len(r.get('pages', []))}")
        lines.append(f"- Images extracted: {len(r.get('images', []))}")
        for p in r.get('pages', [])[:8]:
            tx = (p.get('text', '') or '').replace('\n', ' ')
            lines.append(f"- Page {p['page']} snippet: {tx[:300]}")
    elif r['type'] == 'docx':
        lines.append(f"- Paragraphs extracted: {len(r.get('paragraphs', []))}")
        lines.append(f"- Tables extracted: {len(r.get('tables', []))}")
        lines.append(f"- Images extracted: {len(r.get('images', []))}")
        for t in r.get('paragraphs', [])[:12]:
            lines.append(f"- Para: {t[:300]}")
    elif r['type'] == 'pptx':
        lines.append(f"- Slides extracted: {len(r.get('slides', []))}")
        lines.append(f"- Images extracted: {len(r.get('images', []))}")
        for s in r.get('slides', [])[:10]:
            joined = ' | '.join(s.get('text', [])[:3])
            lines.append(f"- Slide {s['slide']} text: {joined[:300]}")
    elif r['type'] == 'conversion':
        lines.append(f"- Conversion status: {r.get('status')}")
        if 'to' in r:
            lines.append(f"- Converted file: {r['to']}")
    elif r['type'] == 'legacy_ppt_pdflike':
        lines.append(f"- Pages extracted: {len(r.get('pages', []))}")
        lines.append(f"- Page images extracted: {len(r.get('images', []))}")
        for p in r.get('pages', [])[:8]:
            tx = (p.get('text', '') or '').replace('\n', ' ')
            ox = (p.get('ocr_page', '') or '').replace('\n', ' ')
            lines.append(f"- Page {p['page']} text snippet: {tx[:220]}")
            lines.append(f"- Page {p['page']} OCR snippet: {ox[:220]}")

OUT_MD.write_text('\n'.join(lines), encoding='utf-8')
print(f'Wrote {OUT_JSON}')
print(f'Wrote {OUT_MD}')
print(f'Total records: {len(records)}')
