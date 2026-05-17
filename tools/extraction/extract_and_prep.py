import os
import argparse
import fitz  # PyMuPDF
from pptx import Presentation
from PIL import Image
import io

def extract_from_pdf(pdf_path, output_img_dir):
    text_content = []
    doc = fitz.open(pdf_path)
    base_name = os.path.splitext(os.path.basename(pdf_path))[0]
    
    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        text = page.get_text("text")
        if text.strip():
            text_content.append(f"--- {base_name} Page {page_num + 1} ---\n{text}")
            
        # Extract images
        image_list = page.get_images(full=True)
        for img_index, img in enumerate(image_list):
            xref = img[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            image_ext = base_image["ext"]
            
            # Save image
            img_filename = f"{base_name}_p{page_num+1}_img{img_index+1}.{image_ext}"
            img_filepath = os.path.join(output_img_dir, img_filename)
            with open(img_filepath, "wb") as f:
                f.write(image_bytes)
                
    return "\n\n".join(text_content)

def extract_from_pptx(pptx_path, output_img_dir):
    text_content = []
    prs = Presentation(pptx_path)
    base_name = os.path.splitext(os.path.basename(pptx_path))[0]
    
    for slide_num, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
                
            # Extract images
            if shape.shape_type == 13: # 13 represents a Picture
                image = shape.image
                image_bytes = image.blob
                image_ext = image.ext
                img_filename = f"{base_name}_s{slide_num+1}_img.{image_ext}"
                img_filepath = os.path.join(output_img_dir, img_filename)
                with open(img_filepath, "wb") as f:
                    f.write(image_bytes)
                    
        if slide_text:
            text_content.append(f"--- {base_name} Slide {slide_num + 1} ---\n" + "\n".join(slide_text))
            
    return "\n\n".join(text_content)

def main():
    parser = argparse.ArgumentParser(description="Automated Study Material Extractor")
    parser.add_argument("input_folder", help="Path to the folder containing PPTX and PDF files")
    parser.add_argument("--output_folder", help="Path to save the generated outputs", default="./output")
    args = parser.parse_args()
    
    in_dir = args.input_folder
    out_dir = args.output_folder
    img_dir = os.path.join(out_dir, "images")
    
    os.makedirs(img_dir, exist_ok=True)
    
    all_text = []
    
    print(f"Scanning {in_dir}...")
    for filename in os.listdir(in_dir):
        filepath = os.path.join(in_dir, filename)
        if filename.lower().endswith(".pdf"):
            print(f"Extracting from PDF: {filename}")
            all_text.append(extract_from_pdf(filepath, img_dir))
        elif filename.lower().endswith(".pptx"):
            print(f"Extracting from PPTX: {filename}")
            all_text.append(extract_from_pptx(filepath, img_dir))
            
    # Save combined raw text
    raw_text_path = os.path.join(out_dir, "raw_extracted_content.txt")
    with open(raw_text_path, "w", encoding="utf-8") as f:
        f.write("\n\n" + "="*50 + "\n\n".join(all_text))
        
    print(f"\nExtraction Complete!")
    print(f"- Raw Text Saved to: {raw_text_path}")
    print(f"- Extracted Images Saved to: {img_dir}")
    
    prompt_path = os.path.join(out_dir, "LLM_PROMPT.md")
    prompt_content = """# LLM Prompt for Generating Study Material

Copy the following prompt and paste it into your LLM (along with the contents of `raw_extracted_content.txt`) to generate the audit-ready LaTeX code:

---

**System Role**: You are an expert LaTeX typesetter and academic content creator.

**Objective**: Convert the provided raw lecture notes (extracted from PDFs/PPTs) into an audit-ready, high-density LaTeX study guide.

**Formatting Requirements**:
1. **Document Class**: Use `\documentclass[11pt,a4paper]{article}` with `geometry` package (1in margins).
2. **Lists**: Use `enumitem` and set global spacing: `\setlist[itemize]{itemsep=0.6em, topsep=0.5em}`
3. **Structure**: Organize strictly by Modules -> Theory -> Architecture -> Solved Numericals.
4. **Images**: Insert relevant extracted images using `\begin{figure}[H]` from the `float` package. Refer to the image names generated during extraction.
5. **Numericals**: Place all numericals at the *end* of each module in a dedicated `\subsection{Solved Numericals}`. 
6. **Numerical Formatting**: Use the `tcolorbox` package for numericals. Each numerical MUST follow this strict structure:
   - **Given:** (bullet points of known variables)
   - **Objective:** (what to find)
   - **Steps:** (numbered list with LaTeX math equations)
   - **Interpretation:** (a 1-2 sentence explanation of the result)
7. **Completeness**: Do not summarize too heavily. Include all formulas, definitions, and theory points found in the raw text. Fill in obvious theoretical gaps to make the notes self-contained.

Please output ONLY the raw LaTeX code starting with `\documentclass` and ending with `\end{document}`.
"""
    with open(prompt_path, "w", encoding="utf-8") as f:
        f.write(prompt_content)
        
    print(f"- LLM Prompt generated at: {prompt_path}")

if __name__ == "__main__":
    main()
